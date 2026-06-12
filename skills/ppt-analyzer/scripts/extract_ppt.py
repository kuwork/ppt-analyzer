#!/usr/bin/env python3
"""
Extract slides to images + notes + text from PPTX.
Uses Aspose.Slides for images and notes, python-pptx for text extraction.
"""
import sys
import os
import re

import aspose.slides as slides
from pptx import Presentation

def sanitize_filename(name):
    return re.sub(r'[\\/:*?"<>|]', '_', name)

def extract_all(pptx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    img_dir = os.path.join(output_dir, "images")
    os.makedirs(img_dir, exist_ok=True)
    
    base_name = sanitize_filename(os.path.splitext(os.path.basename(pptx_path))[0])
    
    # --- Aspose: images + notes ---
    pres_aspose = slides.Presentation(pptx_path)
    total = len(pres_aspose.slides)
    print(f"[Aspose] Loaded {total} slides")
    
    slide_data = []
    for i in range(total):
        slide = pres_aspose.slides[i]
        
        # Export image
        img_filename = f"slide_{i+1:03d}.png"
        img_path = os.path.join(img_dir, img_filename)
        if not os.path.exists(img_path):
            img = slide.get_image()
            img.save(img_path, slides.ImageFormat.PNG)
        
        # Extract notes
        notes_text = ""
        if slide.notes_slide_manager.notes_slide is not None:
            notes_text = slide.notes_slide_manager.notes_slide.notes_text_frame.text.strip()
        
        slide_data.append({
            "slide_num": i + 1,
            "img_path": img_path,
            "img_rel": os.path.join("images", img_filename).replace("\\", "/"),
            "notes": notes_text
        })
    
    # --- python-pptx: text extraction ---
    pres_pptx = Presentation(pptx_path)
    for i, slide in enumerate(pres_pptx.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            elif shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))
        slide_data[i]["texts"] = texts
    
    # --- Write Markdown ---
    md_path = os.path.join(output_dir, f"{base_name}_content.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(f"# {os.path.basename(pptx_path)}\n\n")
        f.write(f"Total slides: {total}\n\n")
        
        for sd in slide_data:
            f.write(f"## Slide {sd['slide_num']}\n\n")
            f.write(f"![Slide {sd['slide_num']}]({sd['img_rel']})\n\n")
            
            if sd["notes"]:
                f.write(f"**Notes (Speaker):**\n\n{sd['notes']}\n\n")
            
            if sd.get("texts"):
                f.write("**Extracted Text:**\n\n")
                for t in sd["texts"]:
                    f.write(f"- {t}\n")
                f.write("\n")
            
            f.write("---\n\n")
    
    print(f"[Done] Images: {img_dir}")
    print(f"[Done] Markdown: {md_path}")
    return md_path, img_dir

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract_ppt_full.py <input.pptx> <output_dir>")
        sys.exit(1)
    extract_all(sys.argv[1], sys.argv[2])
