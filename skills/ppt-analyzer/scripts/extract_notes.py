#!/usr/bin/env python3
"""
通用 PPT 备注与文本提取脚本。
提取每页的演讲者备注和可见文本内容，输出为 Markdown 格式。
依赖：pip install python-pptx
"""
import sys
import os
from pptx import Presentation


def extract_notes_and_text(pptx_path, output_md_path=None):
    prs = Presentation(pptx_path)
    total = len(prs.slides)
    lines = []

    for i, slide in enumerate(prs.slides, start=1):
        # --- 提取备注 ---
        note_text = ""
        if slide.has_notes_slide:
            note_slide = slide.notes_slide
            text_frame = note_slide.notes_text_frame
            note_text = text_frame.text.strip()

        # --- 提取可见文本（去重、去空） ---
        visible_texts = []
        seen = set()
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and t not in seen:
                        visible_texts.append(t)
                        seen.add(t)
            elif shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        t = " | ".join(row_texts)
                        if t not in seen:
                            visible_texts.append(t)
                            seen.add(t)

        # --- 写入 Markdown ---
        lines.append(f"## 第 {i} 页")
        lines.append("")
        lines.append(f"**备注**：{note_text if note_text else '（无）'}")
        lines.append("")
        if visible_texts:
            lines.append("**页面文本**：")
            lines.append("")
            for t in visible_texts:
                lines.append(f"- {t}")
            lines.append("")
        lines.append("")

    md_content = "\n".join(lines)

    if output_md_path:
        os.makedirs(os.path.dirname(output_md_path) if os.path.dirname(output_md_path) else ".", exist_ok=True)
        with open(output_md_path, "w", encoding="utf-8") as f:
            f.write(md_content)
        print(f"已保存到: {output_md_path}  (共 {total} 页)")
    else:
        print(md_content)
        print(f"\n总共 {total} 页")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_notes.py <pptx路径> [输出md路径]")
        sys.exit(1)
    pptx_path = sys.argv[1]
    output_md = sys.argv[2] if len(sys.argv) > 2 else None
    extract_notes_and_text(pptx_path, output_md)
