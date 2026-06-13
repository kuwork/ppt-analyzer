#!/usr/bin/env python3
"""
PPT 演讲者备注批量写入脚本。
支持单页写入或批量写入，输入为页码范围及文本数组。

用法:
    python write_notes.py <pptx路径> <备注JSON文件>

备注JSON文件格式:
    {
        "pages": [
            {"page": 1, "note": "第一页备注内容"},
            {"page": 3, "note": "第三页备注内容"}
        ]
    }

或批量范围格式:
    {
        "pages": [
            {"page": "1-5", "notes": ["备注1", "备注2", "备注3", "备注4", "备注5"]}
        ]
    }
"""
import sys
import json
import os
from pptx import Presentation


def write_notes(pptx_path, notes_data, output_path=None):
    """
    将备注写入PPT文件。
    
    Args:
        pptx_path: 源PPT文件路径
        notes_data: 备注数据，格式为 {"pages": [{"page": 1, "note": "..."}, ...]}
        output_path: 输出文件路径，默认覆盖原文件或添加后缀
    """
    prs = Presentation(pptx_path)
    total = len(prs.slides)
    
    if output_path is None:
        base, ext = os.path.splitext(pptx_path)
        output_path = f"{base}_已更新备注{ext}"
    
    written = 0
    errors = []
    
    for item in notes_data.get("pages", []):
        # 处理单页
        if "page" in item and "note" in item:
            page_num = item["page"]
            note_text = item["note"]
            
            if isinstance(page_num, int):
                if 1 <= page_num <= total:
                    try:
                        slide = prs.slides[page_num - 1]
                        # 确保有备注幻灯片
                        if not slide.has_notes_slide:
                            # 添加备注幻灯片
                            slide.notes_slide = prs.notes_slide_layout
                        
                        notes_slide = slide.notes_slide
                        notes_text_frame = notes_slide.notes_text_frame
                        notes_text_frame.clear()
                        
                        # 设置文本
                        p = notes_text_frame.paragraphs[0]
                        p.text = note_text
                        
                        written += 1
                        print(f"已写入第 {page_num} 页备注")
                    except Exception as e:
                        errors.append(f"第 {page_num} 页写入失败: {e}")
                else:
                    errors.append(f"页码 {page_num} 超出范围 (1-{total})")
        
        # 处理批量范围
        elif "page" in item and "notes" in item:
            range_str = item["page"]
            notes_list = item["notes"]
            
            if isinstance(range_str, str) and "-" in range_str:
                start, end = map(int, range_str.split("-"))
                if end - start + 1 == len(notes_list):
                    for i, note_text in enumerate(notes_list):
                        page_num = start + i
                        if 1 <= page_num <= total:
                            try:
                                slide = prs.slides[page_num - 1]
                                if not slide.has_notes_slide:
                                    slide.notes_slide = prs.notes_slide_layout
                                
                                notes_slide = slide.notes_slide
                                notes_text_frame = notes_slide.notes_text_frame
                                notes_text_frame.clear()
                                
                                p = notes_text_frame.paragraphs[0]
                                p.text = note_text
                                
                                written += 1
                            except Exception as e:
                                errors.append(f"第 {page_num} 页写入失败: {e}")
                        else:
                            errors.append(f"页码 {page_num} 超出范围 (1-{total})")
                    print(f"已写入第 {start}-{end} 页批量备注")
                else:
                    errors.append(f"范围 {range_str} 与备注数量 {len(notes_list)} 不匹配")
    
    # 保存文件
    prs.save(output_path)
    print(f"\n保存到: {output_path}")
    print(f"成功写入: {written} 页")
    if errors:
        print(f"错误: {len(errors)} 项")
        for err in errors:
            print(f"  - {err}")
    
    return output_path, written, errors


def write_notes_from_markdown(pptx_path, md_path, output_path=None):
    """
    从备注提取Markdown文件读取备注并写入PPT。
    
    Args:
        pptx_path: 源PPT文件路径
        md_path: 备注提取.md文件路径
        output_path: 输出文件路径
    """
    import re
    
    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()
    
    # 解析每页备注
    pages = []
    pattern = r"## 第 (\d+) 页\s*\n\s*\*\*备注\*\*：(.*?)(?=\n\s*\n|\n## 第 |\Z)"
    matches = re.findall(pattern, content, re.DOTALL)
    
    for page_num_str, note_text in matches:
        note_clean = note_text.strip()
        if note_clean and note_clean != "（无）":
            pages.append({
                "page": int(page_num_str),
                "note": note_clean
            })
    
    notes_data = {"pages": pages}
    return write_notes(pptx_path, notes_data, output_path)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法:")
        print("  1. 从JSON文件写入: python write_notes.py <pptx路径> <json路径> [输出路径]")
        print("  2. 从Markdown写入: python write_notes.py --from-md <pptx路径> <md路径> [输出路径]")
        sys.exit(1)
    
    if sys.argv[1] == "--from-md":
        pptx_path = sys.argv[2]
        md_path = sys.argv[3]
        output_path = sys.argv[4] if len(sys.argv) > 4 else None
        write_notes_from_markdown(pptx_path, md_path, output_path)
    else:
        pptx_path = sys.argv[1]
        json_path = sys.argv[2]
        output_path = sys.argv[3] if len(sys.argv) > 3 else None
        
        with open(json_path, "r", encoding="utf-8") as f:
            notes_data = json.load(f)
        
        write_notes(pptx_path, notes_data, output_path)
